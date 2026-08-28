from __future__ import annotations

from io import BytesIO

import pandas as pd

from app.ingestion.extraction_contract import (
    ArtifactExtractionResult,
    ExtractedDataset,
    ExtractedEvidence,
)

# python-pptx (and its lxml dependency) is optional at runtime: some
# deployment environments block lxml's compiled extension outright (e.g. an
# Application Control / EDR policy). That must never prevent the app from
# starting or prevent unrelated CSV/JSON/XLSX artifacts from being
# processed -- this parser stays registered and its class stays
# importable/instantiable either way; only extract() is gated, and it
# reports an honest ArtifactExtractionResult(status="unavailable") instead
# of raising ImportError/NameError at call time.
_IMPORT_ERROR: str | None = None
try:
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover -- environment dependent
    _IMPORT_ERROR = str(exc)


class PowerPointParser:
    code = "pptx"
    version = "1.0"

    def supports(self, mime_type: str, extension: str) -> bool:
        return extension.lower() == ".pptx" or mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def is_available(self) -> bool:
        return _IMPORT_ERROR is None

    def extract(self, raw_bytes: bytes, filename: str) -> ArtifactExtractionResult:
        if _IMPORT_ERROR is not None:
            return ArtifactExtractionResult(
                parser_code=self.code,
                parser_version=self.version,
                status="unavailable",
                warnings=[
                    "PPTX parsing dependency (python-pptx/lxml) is unavailable in this "
                    f"environment: {_IMPORT_ERROR}"
                ],
                extraction_metadata={"reason": "dependency_unavailable"},
            )
        try:
            presentation = Presentation(BytesIO(raw_bytes))
        except Exception as exc:  # noqa: BLE001
            return ArtifactExtractionResult(
                parser_code=self.code,
                parser_version=self.version,
                status="failed",
                warnings=[f"PPTX parse failed: {exc}"],
            )
        evidence = []
        datasets = []
        for slide_index, slide in enumerate(presentation.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    if len(rows) >= 2:
                        dataframe = pd.DataFrame(rows[1:], columns=rows[0])
                        datasets.append(
                            ExtractedDataset(
                                label=f"{filename}:slide_{slide_index + 1}_table",
                                dataframe=dataframe,
                                lineage={"slide_number": slide_index + 1},
                            )
                        )
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text
            slide_text = "\n".join(texts)
            if slide_text or notes:
                evidence.append(
                    ExtractedEvidence(
                        evidence_type="text_block",
                        content="\n".join(filter(None, [slide_text, notes])),
                        lineage={"slide_number": slide_index + 1},
                    )
                )
        return ArtifactExtractionResult(
            parser_code=self.code,
            parser_version=self.version,
            status="extracted" if (evidence or datasets) else "partial",
            evidence_objects=evidence,
            datasets=datasets,
            extraction_metadata={"slide_count": len(presentation.slides)},
        )
